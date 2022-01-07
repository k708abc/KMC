from pptx import Presentation
from pptx.util import Inches, Pt
import os
import datetime


def rec_ppt(
    params,
    minute,
    second,
    img_names,
    hist_names,
    time,
    coverage,
    dir_name,
    time_per_dep,
    growth_mode,
    mode_val,
    other_modes,
):
    ppt_name = dir_name + "Layer_analysis_results.pptx"
    if os.path.exists(ppt_name):
        prs = Presentation(ppt_name)
    else:
        prs = Presentation()
        prs.slide_height = Inches(7.5)
        prs.slide_width = Inches(13.33)
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = "Layer analysis calculation results"
    # First page
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    width = height = Inches(1)
    top = Inches(0)
    left = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Parameters"
    p.font.size = Pt(28)
    shapes = slide.shapes
    left = Inches(0.1)
    top = Inches(1)
    rows = 2
    cols = 9
    width = Inches(13)
    height = Inches(1)
    # record parameters
    table0 = shapes.add_table(rows, cols, left, top, width, height).table
    table0.cell(0, 0).text = "Num. cells"
    table0.cell(0, 1).text = "Z units"
    table0.cell(0, 2).text = "T (K)"
    table0.cell(0, 3).text = "kbT"
    table0.cell(0, 4).text = "dep. rate (ML/min)"
    table0.cell(0, 5).text = "dep. time (min)"
    table0.cell(0, 6).text = "prefactor"
    table0.cell(0, 7).text = "Cal. time (s)"
    table0.cell(0, 8).text = "Finishing time"
    table0.cell(1, 0).text = str(params.cell_size_xy)
    table0.cell(1, 1).text = str(params.cell_size_z)
    table0.cell(1, 2).text = str(params.temperature)
    table0.cell(1, 3).text = str("{:.3g}".format(params.temperature_eV))
    table0.cell(1, 4).text = str(params.dep_rate)
    table0.cell(1, 5).text = str(params.dep_time)
    table0.cell(1, 6).text = str("{:.1E}".format(float(params.prefactor)))
    table0.cell(1, 7).text = str(minute) + " m " + str(second) + " s"
    dt_now = datetime.datetime.now()
    table0.cell(1, 8).text = str(dt_now.strftime("%Y/%m/%d/ %H:%M:%S"))
    #
    left = Inches(0.1)
    top = Inches(2.5)
    rows = 3
    cols = 3
    width = Inches(6)
    height = Inches(1)

    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 0).text = "Put at first"
    table.cell(0, 1).text = "Transformation"
    table.cell(0, 2).text = "time/event (ms)"
    table.cell(1, 0).text = str(params.first_put_check)
    table.cell(1, 1).text = str(params.trans_check)
    table.cell(2, 0).text = str(params.first_put_num)
    table.cell(2, 1).text = str(params.trans_num)
    table.cell(2, 2).text = str(time_per_dep)
    #
    left = Inches(0.1)
    top = Inches(4.8)
    rows = 3
    cols = 5
    width = Inches(13)
    height = Inches(1)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 1).text = "First"
    table.cell(0, 2).text = "Second"
    table.cell(0, 3).text = "Third"
    table.cell(0, 4).text = "Else"
    table.cell(1, 0).text = "Diffusion (eV)"
    table.cell(1, 1).text = str(params.energies_diffusion["first"])
    table.cell(1, 2).text = str(params.energies_diffusion["second"])
    table.cell(1, 3).text = str(params.energies_diffusion["third"])
    table.cell(1, 4).text = str(params.energies_diffusion["upper"])

    #
    table.cell(2, 0).text = "Binding (eV)"
    table.cell(2, 1).text = str(params.energies_binding["first"])
    table.cell(2, 2).text = str(params.energies_binding["second"])
    table.cell(2, 3).text = str(params.energies_binding["third"])
    table.cell(2, 4).text = str(params.energies_binding["upper"])
    ##

    width = height = Inches(1)
    top = Inches(5.9)
    left = Inches(0.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Comment: " + "\n" + params.comments
    p.font.size = Pt(20)

    #
    num_ims = len(img_names)
    imn = 0
    while imn < num_ims:
        slide = prs.slides.add_slide(blank_slide_layout)
        width = height = Inches(1)
        top = Inches(-0.1)
        left = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "Results"
        p.font.size = Pt(28)
        left0 = 0.2
        top0 = 0.7
        height = Inches(2.5)
        height_w = Inches(1)
        #
        for i in range(0, 3):
            for k in range(0, 4):
                if imn >= num_ims:
                    break
                else:
                    left = Inches(left0 + k * 3.2)
                    top = Inches(top0 + i * 2.1)
                    slide.shapes.add_picture(img_names[imn], left, top, height=height)
                    txBox = slide.shapes.add_textbox(
                        left, top - Inches(0.2), width, height_w
                    )
                    tf = txBox.text_frame
                    p = tf.add_paragraph()

                    p.text = (
                        str(int(time[imn]))
                        + " s, "
                        + str("{:.2f}".format(coverage[imn]))
                        + " ML"
                    )
                    p.font.size = Pt(20)

                    imn = imn + 1

    # Third slide
    num_ims = len(hist_names)
    imn = 0
    while imn < num_ims:
        slide = prs.slides.add_slide(blank_slide_layout)
        width = height = Inches(1)
        top = Inches(-0.1)
        left = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "Results: layer analysis"
        p.font.size = Pt(28)
        left0 = 0.2
        top0 = 0.7
        height = Inches(2)
        height_w = Inches(1)
        for i in range(0, 3):
            for k in range(0, 4):
                if imn >= num_ims:
                    break
                else:
                    left = Inches(left0 + k * 3.2)
                    top = Inches(top0 + i * 2.1)
                    slide.shapes.add_picture(
                        hist_names[imn], left, top + Inches(0.2), height=height
                    )

                    txBox = slide.shapes.add_textbox(
                        left, top - Inches(0.2), width, height_w
                    )
                    tf = txBox.text_frame

                    p = tf.add_paragraph()
                    p.text = (
                        str(int(time[imn]))
                        + " s, "
                        + str("{:.2f}".format(coverage[imn]))
                        + " ML"
                    )
                    p.font.size = Pt(20)

                    imn = imn + 1
    #
    slide = prs.slides.add_slide(blank_slide_layout)
    width = height = Inches(1)
    top = Inches(-0.1)
    left = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Growth mode"
    p.font.size = Pt(28)
    n_growth = len(growth_mode)
    #
    left = Inches(0.1)
    top = Inches(2.5)
    rows = 5
    cols = 9
    width = Inches(13)
    height = Inches(1)
    inch = 1
    imn = 0
    while imn < n_growth:
        top = Inches(inch)
        table_g = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table_g.cell(0, 0).text = "Coverage"
        table_g.cell(1, 0).text = "Ag"
        table_g.cell(2, 0).text = "1st layer"
        table_g.cell(3, 0).text = "2nd layer"
        table_g.cell(4, 0).text = "Multi layer"
        for i in range(8):
            if imn >= n_growth:
                pass
            else:
                table_g.cell(0, i + 1).text = str(coverage[imn])
                table_g.cell(1, i + 1).text = str(growth_mode[imn][0])
                table_g.cell(2, i + 1).text = str(growth_mode[imn][1])
                table_g.cell(3, i + 1).text = str(growth_mode[imn][2])
                table_g.cell(4, i + 1).text = str(growth_mode[imn][3])
                imn += 1
        inch += 1.8
    #
    top = Inches(inch)
    rows = 2
    cols = 4
    table_mode = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table_mode.cell(0, 0).text = "Mode"
    table_mode.cell(0, 1).text = "First layer at 1ML"
    table_mode.cell(0, 2).text = "Second layer at 2ML"
    table_mode.cell(0, 3).text = "First layer/100-Multi at 2ML"
    table_mode.cell(1, 0).text = str(mode_val[3])
    table_mode.cell(1, 1).text = str(mode_val[0])
    table_mode.cell(1, 2).text = str(mode_val[1])
    table_mode.cell(1, 3).text = str(mode_val[2])
    #
    #
    slide = prs.slides.add_slide(blank_slide_layout)
    width = height = Inches(1)
    top = Inches(-0.1)
    left = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Other modes"
    p.font.size = Pt(28)
    #
    width = Inches(13)
    height = Inches(1)
    left = Inches(0.1)
    top = Inches(1)
    repeat = 1
    for mode_val in other_modes:
        top += Inches(1)
        rows = 2
        cols = 4
        table_mode2 = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table_mode2.cell(0, 0).text = "Mode" + str(repeat)
        table_mode2.cell(0, 1).text = "value 1"
        table_mode2.cell(0, 2).text = "value 2"
        table_mode2.cell(0, 3).text = "value 3"
        table_mode2.cell(1, 0).text = str(mode_val[3])
        table_mode2.cell(1, 1).text = str(mode_val[0])
        table_mode2.cell(1, 2).text = str(mode_val[1])
        table_mode2.cell(1, 3).text = str(mode_val[2])
        repeat += 1
    #
    #
    #
    #
    #
    slide = prs.slides.add_slide(blank_slide_layout)
    width = height = Inches(1)
    top = Inches(-0.1)
    left = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Number of events per deposition" + "\n" + "Change of coverage"
    p.font.size = Pt(28)
    height = Inches(5.5)
    top = Inches(2)
    left = Inches(0)
    slide.shapes.add_picture(
        dir_name + "Num_events_per_dep.png", left, top, height=height
    )
    left = Inches(6)
    slide.shapes.add_picture(dir_name + "Coverage_change.png", left, top, height=height)
    #

    rec_num = 0
    while rec_num in (0, 1):
        try:
            prs.save(ppt_name)
            rec_num = 2
        except:
            if rec_num == 0:
                print("Close the ppt")
                rec_num = 1
            else:
                pass
