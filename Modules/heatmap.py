import matplotlib.patches as patches
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
import os


def form_heatmap(growth_mode, rec_name, E1, E2, diff1, diff2):
    fig = plt.figure(figsize=(6, 6))
    ax = fig.add_subplot(111)
    for i in range(len(E1)):
        for k in range(len(E2)):
            RGB = [0, 0, 0]
            s1 = growth_mode[i][k][0]
            s2 = growth_mode[i][k][1]

            sn1 = s1 + s2 - 100
            sn2 = -s1 + s2

            if sn1 >= 0:
                RGB[1] = sn1 / 100
            else:
                RGB[2] = abs(sn1 / 100)

            if sn2 <= 0:
                RGB[0] = abs(sn2 / 100)

            r = patches.Rectangle(
                xy=(E1[i] - diff1 / 2, E2[k] - diff2 / 2),
                width=diff1,
                height=diff2,
                color=RGB,
                fill=True,
            )

            ax.add_patch(r)

    ax.set_xlim(min(E1) - diff1 / 2, max(E1) + diff1 / 2)
    ax.set_ylim(min(E2) - diff2 / 2, max(E2) + diff2 / 2)

    E1_2 = []
    E2_2 = []

    if len(E1) >= 10:
        for i in range(len(E1)):
            if i % 2 == 0:
                E1_2.append(E1[i])
    else:
        for i in range(len(E1)):
            E1_2.append(E1[i])

    if len(E2) >= 10:
        for i in range(len(E2)):
            if i % 2 == 0:
                E2_2.append(E2[i])
    else:
        for i in range(len(E2)):
            E2_2.append(E2[i])

    ax.set_xticks(E1_2)
    ax.set_yticks(E2_2)
    ax.set_xlabel("First layer")
    ax.set_ylabel("Multi layer")
    labels = ax.get_xticklabels()
    plt.setp(labels, rotation=45, fontsize=16)
    plt.tick_params(labelsize=16)
    plt.tight_layout()
    #
    dir_name = rec_name + "/"
    heat_name = dir_name + "heatmap.png"
    if os.path.exists(dir_name) is False:
        os.mkdir(dir_name)
    plt.savefig(heat_name)
    #
    heat_text = dir_name + "heat_values.txt"
    file_data = open(heat_text, "w")

    file_data.write("E1" + "\t" + "E2" + "\t" + "ML" + "\t" + "BL" + "\n")

    for i in range(len(E1)):
        for k in range(len(E2)):
            file_data.write(
                str(E1[i])
                + "\t"
                + str(E2[k])
                + "\t"
                + str(growth_mode[i][k][0])
                + "\t"
                + str(growth_mode[i][k][1])
                + "\n"
            )

    file_data.close()
    #
    ppt_name = dir_name + "Layer_analysis_results.pptx"
    if os.path.exists(ppt_name):
        prs = Presentation(ppt_name)
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        width = height = Inches(1)
        top = Inches(0)
        left = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "Heat map"
        height = Inches(5.5)
        top = Inches(2)
        left = Inches(0)
        slide.shapes.add_picture(heat_name, left, top, height=height)
        prs.save(ppt_name)
