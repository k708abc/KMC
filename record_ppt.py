from pptx import Presentation
from pptx.util import Inches, Pt
import os


def rec_ppt(params, minute, second, img_names, hist_names, time, coverage, dir_name):
    ppt_name = dir_name + "Layer_analysis_results.pptx"
    if os.path.exists(ppt_name):
        prs = Presentation(ppt_name)
    else:
        prs = Presentation()
        prs.slide_height = Inches(7.5)
        prs.slide_width = Inches(13.33)
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = "Layer analysis calculation results"
    # First page
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    width = height = Inches(1)
    top = Inches(0)
    left = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Parameters"
    p.font.size = Pt(28)
    shapes = slide.shapes
    left = Inches(0.5)
    top = Inches(1.5)
    rows = 2
    cols = 11
    width = Inches(12)
    height = Inches(1)
    # record parameters
    table0 = shapes.add_table(rows, cols, left, top, width, height).table
    table0.cell(0, 0).text = "Num. cells"
    table0.cell(0, 1).text = "Z units"
    table0.cell(0, 2).text = "T (K)"
    table0.cell(0, 3).text = "kbT"
    table0.cell(0, 4).text = "dep. rate (ML/min)"
    table0.cell(0, 5).text = "dep. time (min)"
    table0.cell(0, 6).text = "post anneal(min)"
    table0.cell(0, 7).text = "prefactor"
    table0.cell(0, 8).text = "transform"
    table0.cell(0, 9).text = "keep defect"
    table0.cell(0, 10).text = "Cal. time (s)"
    table0.cell(1, 0).text = str(params.n_cell_init)
    table0.cell(1, 1).text = str(params.z_unit_init)
    table0.cell(1, 2).text = str(params.temperature)
    table0.cell(1, 3).text = str("{:.3g}".format(params.temperature_eV))
    table0.cell(1, 4).text = str(params.dep_rate)
    table0.cell(1, 5).text = str(params.dep_time)
    table0.cell(1, 6).text = str(params.post_anneal)
    table0.cell(1, 7).text = str("{:.1E}".format(params.prefactor))
    table0.cell(1, 8).text = "pass"
    table0.cell(1, 9).text = "pass"
    """
    if bln_tr.get() == True:
        table0.cell(1, 8).text = "on"
    else:
        table0.cell(1, 8).text = "off"
    if bln_def.get() == True:
        table0.cell(1, 9).text = "on"
    else:
        table0.cell(1, 9).text = "off"
    """
    table0.cell(1, 10).text = str(minute) + " m " + str(second) + " s"
    #
    left = Inches(0.5)
    top = Inches(3.5)
    rows = 2
    cols = 11
    width = Inches(12)
    height = Inches(1)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 1).text = "Ag-Si"
    table.cell(0, 2).text = "Si(1-2)"
    table.cell(0, 3).text = "Si(2-3)"
    table.cell(0, 4).text = "Si(3-4)"
    table.cell(0, 5).text = "Si(4-5)"
    table.cell(0, 6).text = "Si(5-6)"
    table.cell(0, 7).text = "Si(intra)"
    table.cell(0, 8).text = "Si(inter)"
    table.cell(0, 9).text = "Ag(top)"
    table.cell(0, 10).text = "Trans."
    table.cell(1, 0).text = "Energy"
    table.cell(1, 1).text = str(params.binding_energies["AgSi"])
    table.cell(1, 2).text = str(params.binding_energies["Si12"])
    table.cell(1, 3).text = str(params.binding_energies["Si23"])
    table.cell(1, 4).text = str(params.binding_energies["Si34"])
    table.cell(1, 5).text = str(params.binding_energies["Si45"])
    table.cell(1, 6).text = str(params.binding_energies["Si56"])
    table.cell(1, 7).text = str(params.binding_energies["Si_intra"])
    table.cell(1, 8).text = str(params.binding_energies["Si_inter"])
    table.cell(1, 9).text = str(params.binding_energies["Agtop"])
    table.cell(1, 10).text = str(params.transformation)
    #
    """
    left = Inches(0.5)
    top = Inches(6.5)
    rows = 2
    cols = 3
    width = Inches(8)
    height = Inches(0.5)

    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 0).text = "Total energy"
    table.cell(0, 1).text = "1 ML"
    table.cell(0, 2).text = "Final"

    table.cell(1, 0).text = str('{:.3g}'.format(E))
    table.cell(1, 1).text = str('{:.3g}'.format(ML_check))
    table.cell(1, 2).text = str('{:.3g}'.format(final_check))
    """
    #
    num_ims = len(img_names)
    imn = 0
    while imn < num_ims:
        slide = prs.slides.add_slide(blank_slide_layout)
        width = height = Inches(1)
        top = Inches(-0.1)
        left = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "Results"
        p.font.size = Pt(28)
        left0 = 0.2
        top0 = 0.7
        height = Inches(2.5)
        height_w = Inches(1)
        #
        for i in range(0, 3):
            for k in range(0, 4):
                if imn >= num_ims:
                    break
                else:
                    left = Inches(left0 + k * 3.2)
                    top = Inches(top0 + i * 2.1)
                    slide.shapes.add_picture(img_names[imn], left, top, height=height)
                    txBox = slide.shapes.add_textbox(
                        left, top - Inches(0.2), width, height_w
                    )
                    tf = txBox.text_frame
                    p = tf.add_paragraph()

                    p.text = (
                        str(int(time[imn]))
                        + " s, "
                        + str("{:.2f}".format(coverage[imn]))
                        + " ML"
                    )
                    p.font.size = Pt(20)

                    imn = imn + 1

    # Third slide
    num_ims = len(hist_names)
    imn = 0
    while imn < num_ims:
        slide = prs.slides.add_slide(blank_slide_layout)
        width = height = Inches(1)
        top = Inches(-0.1)
        left = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "Results: layer analysis"
        p.font.size = Pt(28)
        left0 = 0.2
        top0 = 0.7
        height = Inches(2)
        height_w = Inches(1)
        for i in range(0, 3):
            for k in range(0, 4):
                if imn >= num_ims:
                    break
                else:
                    left = Inches(left0 + k * 3.2)
                    top = Inches(top0 + i * 2.1)
                    slide.shapes.add_picture(
                        hist_names[imn], left, top + Inches(0.2), height=height
                    )

                    txBox = slide.shapes.add_textbox(
                        left, top - Inches(0.2), width, height_w
                    )
                    tf = txBox.text_frame

                    p = tf.add_paragraph()
                    p.text = (
                        str(int(time[imn]))
                        + " s, "
                        + str("{:.2f}".format(coverage[imn]))
                        + " ML"
                    )
                    p.font.size = Pt(20)

                    imn = imn + 1
    """
    slide = prs.slides.add_slide(blank_slide_layout)

    width = height = Inches(1)
    top = Inches(-0.1)
    left = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    p = tf.add_paragraph()
    p.text = "Results: Coverage dependence"
    p.font.size = Pt(28)

    height = Inches(5.5)
    top = Inches(1)
    left = Inches(0.7)

    file_name = entry_rec.get() + "_coverage" + ".png"

    slide.shapes.add_picture(file_name, left, top, height=height)
    """
    prs.save(ppt_name)
