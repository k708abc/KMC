#!/usr/bin/env python3

import copy
import math
import time
from typing import List, Dict
import tkinter as tk
import os
import random
import matplotlib.pyplot as plt
import matplotlib.patches as pat
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
from PIL import Image, ImageGrab
from pptx import Presentation
from pptx.util import Inches, Pt
from preference_window import Window

input_params: Dict[str, float] = {}
lattice: Dict[str, list] = {}
atom_set: Dict[str, int] = {}
bonds: Dict[str, list] = {}
event: Dict[str, list] = {}
event_time: Dict[str, list] = {}
event_time_tot: Dict[str, float] = {}


def show_current():
    global atom_set, lattice, c_num, max_layer
    print("current saving")
    # image
    fig = plt.figure()
    ax = fig.add_subplot(111)
    length = nl
    p = pat.Polygon(
        xy=[
            (0, 0),
            (unit_x[0] * length, unit_x[1] * length),
            ((unit_x[0] + unit_y[0]) * length, (unit_x[1] + unit_y[1]) * length),
            (unit_y[0] * length, unit_y[1] * length),
        ],
        fc=(0.5, 0.5, 0.5),
        ec=(0, 0, 0),
    )
    ax.add_patch(p)
    atom_set_n = atom_set
    for z in range(0, max_layer):
        for i in range(len(atom_set_n)):
            for k in range(len(atom_set_n[i])):
                if atom_set_n[i][k][z] == 0:
                    pass
                else:
                    xp = lattice[i][k][z][0]
                    yp = lattice[i][k][z][1]
                    zp = lattice[i][k][z][2]

                    if z % 2 == 0:
                        t1x = (
                            xp * unit_x[0]
                            + yp * unit_y[0]
                            - (unit_x[0] + unit_y[0]) / 3
                        )
                        t1y = (
                            xp * unit_x[1]
                            + yp * unit_y[1]
                            - (unit_x[1] + unit_y[1]) / 3
                        )
                        t2x = t1x + unit_x[0]
                        t2y = t1y + unit_x[1]
                        t3x = t1x + unit_y[0]
                        t3y = t1y + unit_y[1]
                    else:
                        t1x = (
                            xp * unit_x[0]
                            + yp * unit_y[0]
                            + (unit_x[0] + unit_y[0]) / 3
                        )
                        t1y = (
                            xp * unit_x[1]
                            + yp * unit_y[1]
                            + (unit_x[1] + unit_y[1]) / 3
                        )

                        t2x = t1x - unit_x[0]
                        t2y = t1y - unit_x[1]

                        t3x = t1x - unit_y[0]
                        t3y = t1y - unit_y[1]

                    color_num = math.floor(z / 2) * 2
                    if z == 1 or z == 0:
                        color = [0, 1, 0]
                    else:
                        color = [color_num / max_layer, 0, 1 - color_num / max_layer]

                    p = pat.Polygon(
                        xy=[(t1x, t1y), (t2x, t2y), (t3x, t3y)], fc=color, ec=color
                    )

                    ax.add_patch(p)

                    if i == 0:
                        if z % 2 == 0:
                            t1x = (
                                xp * unit_x[0]
                                + yp * unit_y[0]
                                - (unit_x[0] + unit_y[0]) / 3
                                + unit_x[0] * length
                            )
                            t1y = (
                                xp * unit_x[1]
                                + yp * unit_y[1]
                                - (unit_x[1] + unit_y[1]) / 3
                                + unit_x[1] * length
                            )

                            t2x = t1x + unit_x[0]
                            t2y = t1y + unit_x[1]

                            t3x = t1x + unit_y[0]
                            t3y = t1y + unit_y[1]
                        else:
                            t1x = (
                                xp * unit_x[0]
                                + yp * unit_y[0]
                                + (unit_x[0] + unit_y[0]) / 3
                                + unit_x[0] * length
                            )
                            t1y = (
                                xp * unit_x[1]
                                + yp * unit_y[1]
                                + (unit_x[1] + unit_y[1]) / 3
                                + unit_x[1] * length
                            )

                            t2x = t1x - unit_x[0]
                            t2y = t1y - unit_x[1]

                            t3x = t1x - unit_y[0]
                            t3y = t1y - unit_y[1]

                        p = pat.Polygon(
                            xy=[(t1x, t1y), (t2x, t2y), (t3x, t3y)], fc=color, ec=color
                        )
                        ax.add_patch(p)

                    if k == 0:
                        if z % 2 == 0:
                            t1x = (
                                xp * unit_x[0]
                                + yp * unit_y[0]
                                - (unit_x[0] + unit_y[0]) / 3
                                + unit_y[0] * length
                            )
                            t1y = (
                                xp * unit_x[1]
                                + yp * unit_y[1]
                                - (unit_x[1] + unit_y[1]) / 3
                                + unit_y[1] * length
                            )

                            t2x = t1x + unit_x[0]
                            t2y = t1y + unit_x[1]

                            t3x = t1x + unit_y[0]
                            t3y = t1y + unit_y[1]
                        else:
                            t1x = (
                                xp * unit_x[0]
                                + yp * unit_y[0]
                                + (unit_x[0] + unit_y[0]) / 3
                                + unit_y[0] * length
                            )
                            t1y = (
                                xp * unit_x[1]
                                + yp * unit_y[1]
                                + (unit_x[1] + unit_y[1]) / 3
                                + unit_y[1] * length
                            )

                            t2x = t1x - unit_x[0]
                            t2y = t1y - unit_x[1]

                            t3x = t1x - unit_y[0]
                            t3y = t1y - unit_y[1]

                        p = pat.Polygon(
                            xy=[(t1x, t1y), (t2x, t2y), (t3x, t3y)], fc=color, ec=color
                        )
                        ax.add_patch(p)

                    if (i == 0) and (k == 0):

                        if z % 2 == 0:
                            t1x = (
                                xp * unit_x[0]
                                + yp * unit_y[0]
                                - (unit_x[0] + unit_y[0]) / 3
                                + unit_y[0] * length
                                + unit_x[0] * length
                            )
                            t1y = (
                                xp * unit_x[1]
                                + yp * unit_y[1]
                                - (unit_x[1] + unit_y[1]) / 3
                                + unit_y[1] * length
                                + unit_x[1] * length
                            )

                            t2x = t1x + unit_x[0]
                            t2y = t1y + unit_x[1]

                            t3x = t1x + unit_y[0]
                            t3y = t1y + unit_y[1]
                        else:
                            t1x = (
                                xp * unit_x[0]
                                + yp * unit_y[0]
                                + (unit_x[0] + unit_y[0]) / 3
                                + unit_y[0] * length
                                + unit_x[0] * length
                            )
                            t1y = (
                                xp * unit_x[1]
                                + yp * unit_y[1]
                                + (unit_x[1] + unit_y[1]) / 3
                                + unit_y[1] * length
                                + unit_x[1] * length
                            )

                            t2x = t1x - unit_x[0]
                            t2y = t1y - unit_x[1]

                            t3x = t1x - unit_y[0]
                            t3y = t1y - unit_y[1]

                        p = pat.Polygon(
                            xy=[(t1x, t1y), (t2x, t2y), (t3x, t3y)], fc=color, ec=color
                        )
                        ax.add_patch(p)
    p = pat.Polygon(
        xy=[(0, 0), (unit_y[0] * length, unit_y[1] * length), (0, unit_y[1] * length)],
        fc=(1, 1, 1),
        ec=(0, 0, 0),
    )
    ax.add_patch(p)
    p = pat.Polygon(
        xy=[
            (unit_x[0] * length, 0),
            ((unit_x[0] + unit_y[0]) * length, 0),
            ((unit_x[0] + unit_y[0]) * length, (unit_x[1] + unit_y[1]) * length),
        ],
        fc=(1, 1, 1),
        ec=(0, 0, 0),
    )
    ax.add_patch(p)
    ax.set_xlim(0, (unit_x[0] + unit_y[0]) * length)
    ax.set_ylim(0, (unit_x[1] + unit_y[1]) * length)
    ax.set_aspect("equal")
    file_name = entry_rec.get() + "_current_" + str(c_num) + ".png"
    fig.savefig(file_name)

    # poscar

    xp = []
    yp = []
    zp = []
    num = 0
    s = atom_set

    for i in range(len(s)):
        for k in range(len(s[i])):
            for z in range(len(s[i][k])):
                if s[i][k][z] != 0:
                    xp.append(lattice[i][k][z][0] / nl)
                    yp.append(lattice[i][k][z][1] / nl)
                    zp.append(lattice[i][k][z][2] / zl / 2.448)
                    num = num + 1

    file_name = entry_rec.get() + "_" + "current" + str(c_num) + ".vasp"
    file_data = open(file_name, "w")
    file_data.write(entry_rec.get() + "\n")
    file_data.write("10" + "\n")
    file_data.write(str(nl) + "\t" + "0" + "\t" + "0" + "\n")
    file_data.write(str(nl / 2) + "\t" + str(nl / 2 * 1.732) + "\t" + "0" + "\n")
    file_data.write("0" + "\t" + "0" + "\t" + str(zl * 2.448) + "\n")
    file_data.write("Si" + "\n")
    file_data.write(str(num) + "\n")
    file_data.write("direct" + "\n")
    for i in range(len(xp)):
        file_data.write(str(xp[i]) + "\t" + str(yp[i]) + "\t" + str(zp[i]) + "\n")
    file_data.close()
    c_num = c_num + 1


def time_progress():
    global rel_time, tot, pbval, rec, rec1
    r2 = 0
    while r2 == 0:
        r2 = random.random()
    ratio = int(rel_time / ttime * 100)
    dt = -math.log(r2) / tot
    rel_time = rel_time + dt
    # recording
    if ratio >= rec:
        rec_atoms()
        rec = rec + rec1


def rec_atoms():
    global a_set_rec, cov_rec, t_rec, cov
    a_set_rec.append([])

    a_set_rec[-1] = copy.deepcopy(atom_set)

    cov_rec.append(cov)
    t_rec.append(rel_time)


def event_progress(i, k, z, j):
    global atom_set, mod, max_layer

    if (
        (i >= len(event))
        or (k >= len(event[i]))
        or (z >= len(event[i][k]))
        or (j >= len(event[i][k][z]))
    ):
        print("Over")
        print("i = " + str(i))
        print("len(i) = " + str(len(event)))
        print("k = " + str(k))
        print("len(k) = " + str(len(event[i])))
        print("z = " + str(z))
        print("len(z) = " + str(len(event[i][k])))
        print("j = " + str(j))
        print("len(j) = " + str(len(event[i][k][z])))
        print("atomset = " + str(atom_set[i][k][z]))
        print("time = " + str(event_tot[i][k][z]))

    s = event[i][k][z][j]
    mod = []
    if atom_set[i][k][z] == 0:
        print("some mistake")

    if s[0] == "m":
        atom_set[i][k][z] = 0
        atom_set[s[1]][s[2]][s[3]] = 2
    elif s[0] == "u":
        atom_set[i][k][z] = 0
        atom_set[s[1]][s[2]][s[3]] = 2
    elif s[0] == "d":
        atom_set[i][k][z] = 0
        atom_set[s[1]][s[2]][s[3]] = 2
    elif s[0] == "t":
        pass
    if s[3] > max_layer:
        max_layer = s[3]

    mod.append([i, k, z])
    mod.append([s[1], s[2], s[3]])

    if bonds[i][k][z] == []:
        pass
    else:
        for a in bonds[i][k][z]:
            if atom_set[a[0]][a[1]][a[2]] != 0:
                mod.append([a[0], a[1], a[2]])

    if bonds[s[1]][s[2]][s[3]] == []:
        pass
    else:
        for a in bonds[s[1]][s[2]][s[3]]:
            if atom_set[a[0]][a[1]][a[2]] != 0:
                mod.append([a[0], a[1], a[2]])


def kMC():
    global tot, d_rate, rel_time, event_num
    event_num = 1
    while rel_time <= dep_time:
        update_progress()

        event_num = event_num + 1

        r1 = random.random()
        rt = r1 * tot
        # judge deposition
        if rt <= d_rate:
            # deposition happens
            deposition()
            # time progess
            time_progress()
            # update events and rates
            update_events()

        else:
            ie = 100
            rt = rt - d_rate
            for i in range(len(event_tot)):
                for k in range(len(event_tot[i])):
                    for z in range(len(event_tot[i][k])):
                        if rt <= 0:
                            pass
                        elif event_tot[i][k][z] == 0:
                            pass
                        elif rt <= event_tot[i][k][z]:
                            ie = i
                            ke = k
                            ze = z
                            trec = rt
                            rt = rt - event_tot[i][k][z]

                        else:
                            rt = rt - event_tot[i][k][z]

            if rt > 0:
                print("strange")
                print("tot = " + str(tot))
                print("remain = " + str(rt))
                print("ie =" + str(ie))
            else:
                for j in range(len(event_time[ie][ke][ze])):
                    # print(event_time[ie][ke][ze][j])
                    if trec <= 0:
                        pass
                    elif event_time[ie][ke][ze][j] == 0:
                        pass
                    elif trec <= event_time[ie][ke][ze][j]:
                        je = j
                        trec = trec - event_time[ie][ke][ze][j]
                    else:
                        trec = trec - event_time[ie][ke][ze][j]

                if trec > 0:
                    print("strange")
                    print("tot = " + str(tot))
                    print("remain = " + str(trec))

                else:

                    # event[ie][ke][ze][je] happens
                    event_progress(ie, ke, ze, je)
                    # time progess
                    time_progress()
                    # update events and rates
                    update_events()

    # deposition finish
    # post annealing
    tot = tot - d_rate

    # repetition untill post annealing finish
    while rel_time <= ttime:
        r1 = random.random()
        rt = r1 * tot

        for i in range(len(event_tot)):
            for k in range(len(event_tot[i])):
                for z in range(len(event_tot[i][k])):
                    if rt <= 0:
                        pass
                    elif event_tot[i][k][z] == 0:
                        pass
                    elif rt <= event_tot[i][k][z]:
                        ie = i
                        ke = k
                        ze = z
                        trec = rt
                        rt = rt - event_tot[i][k][z]
                    else:
                        rt = rt - event_tot[i][k][z]

        if rt > 0:
            print("strange")
            print("tot = " + str(tot))
            print("remain = " + str(rt))
            print("ie =" + str(ie))

        else:

            for j in range(len(event_time[ie][ke][ze])):
                if trec <= 0:
                    pass
                elif event_time[ie][ke][ze] == 0:
                    pass
                elif trec <= event_time[ie][ke][ze][j]:
                    je = j
                    trec = trec - event_time[ie][ke][ze][j]

            if trec > 0:
                pass
            else:
                # event[ie][ke][ze][je] happens
                event_progress(ie, ke, ze, je)
                # time progess
                time_progress()
                # update events and rates
                update_events()


def rec_pos():
    # record the position of atoms in poscar form
    for n in range(len(a_set_rec)):
        xp = []
        yp = []
        zp = []
        num = 0
        s = a_set_rec[n]

        for i in range(len(s)):
            for k in range(len(s[i])):
                for z in range(len(s[i][k])):
                    if s[i][k][z] != 0:
                        xp.append(lattice[i][k][z][0] / nl)
                        yp.append(lattice[i][k][z][1] / nl)
                        zp.append(lattice[i][k][z][2] / zl / 2.448)
                        num = num + 1

        file_name: str = (
            entry_rec.get() + "_" + str(int(t_rec[n])) + "s" + "_" + str(n) + ".vasp"
        )
        file_data = open(file_name, "w")
        file_data.write(entry_rec.get() + "\n")
        file_data.write("10" + "\n")

        file_data.write(str(nl) + "\t" + "0" + "\t" + "0" + "\n")
        file_data.write(str(nl / 2) + "\t" + str(nl / 2 * 1.732) + "\t" + "0" + "\n")
        file_data.write("0" + "\t" + "0" + "\t" + str(zl * 2.448) + "\n")
        file_data.write("Si" + "\n")
        file_data.write(str(num) + "\n")
        file_data.write("direct" + "\n")
        for i in range(len(xp)):
            file_data.write(str(xp[i]) + "\t" + str(yp[i]) + "\t" + str(zp[i]) + "\n")
        file_data.close()


def figure_draw_rec():
    for n in range(len(a_set_rec)):
        fig = plt.figure()
        ax = fig.add_subplot(111)
        kx = figure_formation(ax, n)
        file_name = (
            entry_rec.get()
            + "_middle_"
            + str(int(t_rec[n]))
            + "s"
            + "_"
            + str(n)
            + ".png"
        )
        fig.savefig(file_name)
        images.append(file_name)


def coverage_color():
    global images, gray, green
    number = 0
    distribution = []
    for i in images:
        distribution.append([0, 0])
        number = number + 1
        img = Image.open(i)
        img.convert("RGB")
        psize = img.size
        for s in range(psize[0]):
            for t in range(psize[1]):
                r, g, b, no = img.getpixel((s, t))
                col = [r, g, b]

                if col == [128, 128, 128]:
                    distribution[-1][0] += 1
                elif col == [0, 255, 0]:
                    distribution[-1][1] += 1
    gray = []
    green = []
    for i in distribution:
        gray.append(i[0])
        green.append(i[1])


def hist_rec():
    global images_h
    images_h = []
    for k in range(len(layers_each)):
        fig = plt.figure()
        ax = fig.add_subplot(111)

        kx = hist_formation(ax, k)

        file_name = (
            entry_rec.get()
            + "_hist_"
            + str(int(t_rec[k]))
            + "s"
            + "_"
            + str(k)
            + ".png"
        )
        fig.savefig(file_name)
        images_h.append(file_name)


def ppt_form():
    global images, images_h
    if os.path.exists("Layer_analysis_results.pptx"):
        prs = Presentation("Layer_analysis_results.pptx")
    else:
        prs = Presentation()
        prs.slide_height = Inches(7.5)
        prs.slide_width = Inches(13.33)
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = "Layer analysis calculation results"
    # First page
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    width = height = Inches(1)
    top = Inches(0)
    left = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Parameters"
    p.font.size = Pt(28)
    shapes = slide.shapes
    left = Inches(0.5)
    top = Inches(1.5)
    rows = 2
    cols = 11
    width = Inches(12)
    height = Inches(1)
    # record parameters
    table0 = shapes.add_table(rows, cols, left, top, width, height).table
    table0.cell(0, 0).text = "Num. cells"
    table0.cell(0, 1).text = "Z units"
    table0.cell(0, 2).text = "T (K)"
    table0.cell(0, 3).text = "kbT"
    table0.cell(0, 4).text = "dep. rate (ML/min)"
    table0.cell(0, 5).text = "dep. time (min)"
    table0.cell(0, 6).text = "post anneal(min)"
    table0.cell(0, 7).text = "prefactor"
    table0.cell(0, 8).text = "transform"
    table0.cell(0, 9).text = "keep defect"
    table0.cell(0, 10).text = "Cal. time (s)"
    table0.cell(1, 0).text = str(nl)
    table0.cell(1, 1).text = str(zl)
    table0.cell(1, 2).text = entry_kbT.get()
    table0.cell(1, 3).text = str("{:.3g}".format(kbt))
    table0.cell(1, 4).text = entry_rate.get()
    table0.cell(1, 5).text = entry_time.get()
    table0.cell(1, 6).text = entry_post.get()
    table0.cell(1, 7).text = entry_pre.get()
    if bln_tr.get() == True:
        table0.cell(1, 8).text = "on"
    else:
        table0.cell(1, 8).text = "off"
    if bln_def.get() == True:
        table0.cell(1, 9).text = "on"
    else:
        table0.cell(1, 9).text = "off"
    table0.cell(1, 10).text = str(minute) + " m " + str(second) + " s"
    left = Inches(0.5)
    top = Inches(3.5)
    rows = 2
    cols = 11
    width = Inches(12)
    height = Inches(1)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 1).text = "Ag-Si"
    table.cell(0, 2).text = "Si(1-2)"
    table.cell(0, 3).text = "Si(2-3)"
    table.cell(0, 4).text = "Si(3-4)"
    table.cell(0, 5).text = "Si(4-5)"
    table.cell(0, 6).text = "Si(5-6)"
    table.cell(0, 7).text = "Si(intra)"
    table.cell(0, 8).text = "Si(inter)"
    table.cell(0, 9).text = "Ag(top)"
    table.cell(0, 10).text = "Trans."
    table.cell(1, 0).text = "Energy"
    table.cell(1, 1).text = entry_AgSi.get()
    table.cell(1, 2).text = entry_Si12.get()
    table.cell(1, 3).text = entry_Si23.get()
    table.cell(1, 4).text = entry_Si34.get()
    table.cell(1, 5).text = entry_Si45.get()
    table.cell(1, 6).text = entry_Si56.get()
    table.cell(1, 7).text = entry_Siel.get()
    table.cell(1, 8).text = entry_Sielin.get()
    table.cell(1, 9).text = entry_Agtp.get()
    table.cell(1, 10).text = entry_tr.get()
    """
    left = Inches(0.5)
    top = Inches(6.5)
    rows = 2
    cols = 3
    width = Inches(8)
    height = Inches(0.5)

    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 0).text = "Total energy"
    table.cell(0, 1).text = "1 ML"
    table.cell(0, 2).text = "Final"

    table.cell(1, 0).text = str('{:.3g}'.format(E))
    table.cell(1, 1).text = str('{:.3g}'.format(ML_check))
    table.cell(1, 2).text = str('{:.3g}'.format(final_check))
    """

    width = height = Inches(1)
    top = Inches(6)
    left = Inches(0.5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    p = tf.add_paragraph()
    p.text = entry_text.get()
    p.font.size = Pt(20)

    # second slide
    slide = prs.slides.add_slide(blank_slide_layout)

    width = height = Inches(1)
    top = Inches(-0.1)
    left = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    p = tf.add_paragraph()
    p.text = "Results"
    p.font.size = Pt(28)

    # put images

    left0 = 0.2
    top0 = 0.7
    height = Inches(2.5)
    height_w = Inches(1)

    num_ims = len(images)
    imn = 0
    slides = 0

    fin = 0

    while fin == 0:
        slides = slides + 1
        if slides != 1:
            slide = prs.slides.add_slide(blank_slide_layout)

            width = height = Inches(1)
            top = Inches(-0.1)
            left = Inches(0.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame

            p = tf.add_paragraph()
            p.text = "Results"
            p.font.size = Pt(28)

            left0 = 0.2
            top0 = 0.7
            height = Inches(2.5)
            height_w = Inches(1)

        for i in range(0, 3):
            for k in range(0, 4):
                if imn >= num_ims:
                    fin = 1
                else:
                    left = Inches(left0 + k * 3.2)
                    top = Inches(top0 + i * 2.1)
                    slide.shapes.add_picture(images[imn], left, top, height=height)

                    txBox = slide.shapes.add_textbox(
                        left, top - Inches(0.2), width, height_w
                    )
                    tf = txBox.text_frame

                    p = tf.add_paragraph()
                    p.text = (
                        str(int(t_rec[imn]))
                        + " s, "
                        + str("{:.2f}".format(cov_rec[imn]))
                        + " ML"
                    )
                    p.font.size = Pt(20)

                    imn = imn + 1
        if imn == num_ims:
            fin = 1

    # Third slide
    slide = prs.slides.add_slide(blank_slide_layout)

    width = height = Inches(1)
    top = Inches(-0.1)
    left = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    p = tf.add_paragraph()
    p.text = "Results: layer analysis"
    p.font.size = Pt(28)

    # put images

    left0 = 0.2
    top0 = 0.7
    height = Inches(2)
    height_w = Inches(1)

    num_ims = len(images_h)
    imn = 0
    slides = 0

    fin = 0

    while fin == 0:
        slides = slides + 1
        if slides != 1:
            slide = prs.slides.add_slide(blank_slide_layout)
            width = height = Inches(1)
            top = Inches(-0.1)
            left = Inches(0.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = "Results: layer analysis"
            p.font.size = Pt(28)
            left0 = 0.2
            top0 = 0.7
            height = Inches(2)
            height_w = Inches(1)
        for i in range(0, 3):
            for k in range(0, 4):
                if imn >= num_ims:
                    fin = 1
                else:
                    left = Inches(left0 + k * 3.2)
                    top = Inches(top0 + i * 2.1)
                    slide.shapes.add_picture(
                        images_h[imn], left, top + Inches(0.2), height=height
                    )

                    txBox = slide.shapes.add_textbox(
                        left, top - Inches(0.2), width, height_w
                    )
                    tf = txBox.text_frame

                    p = tf.add_paragraph()
                    p.text = (
                        str(int(t_rec[imn]))
                        + " s, "
                        + str("{:.2f}".format(cov_rec[imn]))
                        + " ML"
                    )
                    p.font.size = Pt(20)

                    imn = imn + 1
        if imn == num_ims:
            fin = 1
    """
    slide = prs.slides.add_slide(blank_slide_layout)

    width = height = Inches(1)
    top = Inches(-0.1)
    left = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    p = tf.add_paragraph()
    p.text = "Results: Coverage dependence"
    p.font.size = Pt(28)

    height = Inches(5.5)
    top = Inches(1)
    left = Inches(0.7)

    file_name = entry_rec.get() + "_coverage" + ".png"

    slide.shapes.add_picture(file_name, left, top, height=height)
    """
    prs.save("Layer_analysis_results.pptx")


def layer_sum():
    global layers_each, max_layer
    layers_each = []
    max_layer = 0
    for s in range(len(a_set_rec)):
        layers_each.append([])
        co = a_set_rec[s]
        for z in range(len(co[0][0])):
            c_each = 0
            if z % 2 == 0:
                count = 0
            for i in range(len(co)):
                for k in range(len(co[i])):
                    if co[i][k][z] != 0:
                        count += 1
                        c_each += 1
            if z % 2 == 1:
                layers_each[-1].append(count)
            if (c_each != 0) and (z >= max_layer):
                max_layer = z


def figure_formation(ax, n):
    length = nl
    p = pat.Polygon(
        xy=[
            (0, 0),
            (unit_x[0] * length, unit_x[1] * length),
            ((unit_x[0] + unit_y[0]) * length, (unit_x[1] + unit_y[1]) * length),
            (unit_y[0] * length, unit_y[1] * length),
        ],
        fc=(0.5, 0.5, 0.5),
        ec=(0, 0, 0),
    )
    ax.add_patch(p)
    atom_set_n = a_set_rec[n]
    for z in range(0, max_layer):
        for i in range(len(atom_set_n)):
            for k in range(len(atom_set_n[i])):
                if atom_set_n[i][k][z] == 0:
                    pass
                else:
                    xp = lattice[i][k][z][0]
                    yp = lattice[i][k][z][1]
                    zp = lattice[i][k][z][2]
                    if z % 2 == 0:
                        t1x = (
                            xp * unit_x[0]
                            + yp * unit_y[0]
                            - (unit_x[0] + unit_y[0]) / 3
                        )
                        t1y = (
                            xp * unit_x[1]
                            + yp * unit_y[1]
                            - (unit_x[1] + unit_y[1]) / 3
                        )

                        t2x = t1x + unit_x[0]
                        t2y = t1y + unit_x[1]

                        t3x = t1x + unit_y[0]
                        t3y = t1y + unit_y[1]

                    else:
                        t1x = (
                            xp * unit_x[0]
                            + yp * unit_y[0]
                            + (unit_x[0] + unit_y[0]) / 3
                        )
                        t1y = (
                            xp * unit_x[1]
                            + yp * unit_y[1]
                            + (unit_x[1] + unit_y[1]) / 3
                        )

                        t2x = t1x - unit_x[0]
                        t2y = t1y - unit_x[1]

                        t3x = t1x - unit_y[0]
                        t3y = t1y - unit_y[1]

                    color_num = math.floor(z / 2) * 2
                    if z == 1 or z == 0:
                        color = [0, 1, 0]
                    else:
                        color = [color_num / max_layer, 0, 1 - color_num / max_layer]

                    p = pat.Polygon(
                        xy=[(t1x, t1y), (t2x, t2y), (t3x, t3y)], fc=color, ec=color
                    )

                    ax.add_patch(p)

                    if i == 0:
                        if z % 2 == 0:
                            t1x = (
                                xp * unit_x[0]
                                + yp * unit_y[0]
                                - (unit_x[0] + unit_y[0]) / 3
                                + unit_x[0] * length
                            )
                            t1y = (
                                xp * unit_x[1]
                                + yp * unit_y[1]
                                - (unit_x[1] + unit_y[1]) / 3
                                + unit_x[1] * length
                            )

                            t2x = t1x + unit_x[0]
                            t2y = t1y + unit_x[1]

                            t3x = t1x + unit_y[0]
                            t3y = t1y + unit_y[1]
                        else:
                            t1x = (
                                xp * unit_x[0]
                                + yp * unit_y[0]
                                + (unit_x[0] + unit_y[0]) / 3
                                + unit_x[0] * length
                            )
                            t1y = (
                                xp * unit_x[1]
                                + yp * unit_y[1]
                                + (unit_x[1] + unit_y[1]) / 3
                                + unit_x[1] * length
                            )

                            t2x = t1x - unit_x[0]
                            t2y = t1y - unit_x[1]

                            t3x = t1x - unit_y[0]
                            t3y = t1y - unit_y[1]

                        p = pat.Polygon(
                            xy=[(t1x, t1y), (t2x, t2y), (t3x, t3y)], fc=color, ec=color
                        )
                        ax.add_patch(p)

                    if k == 0:
                        if z % 2 == 0:
                            t1x = (
                                xp * unit_x[0]
                                + yp * unit_y[0]
                                - (unit_x[0] + unit_y[0]) / 3
                                + unit_y[0] * length
                            )
                            t1y = (
                                xp * unit_x[1]
                                + yp * unit_y[1]
                                - (unit_x[1] + unit_y[1]) / 3
                                + unit_y[1] * length
                            )

                            t2x = t1x + unit_x[0]
                            t2y = t1y + unit_x[1]

                            t3x = t1x + unit_y[0]
                            t3y = t1y + unit_y[1]
                        else:
                            t1x = (
                                xp * unit_x[0]
                                + yp * unit_y[0]
                                + (unit_x[0] + unit_y[0]) / 3
                                + unit_y[0] * length
                            )
                            t1y = (
                                xp * unit_x[1]
                                + yp * unit_y[1]
                                + (unit_x[1] + unit_y[1]) / 3
                                + unit_y[1] * length
                            )

                            t2x = t1x - unit_x[0]
                            t2y = t1y - unit_x[1]

                            t3x = t1x - unit_y[0]
                            t3y = t1y - unit_y[1]

                        p = pat.Polygon(
                            xy=[(t1x, t1y), (t2x, t2y), (t3x, t3y)], fc=color, ec=color
                        )
                        ax.add_patch(p)

                    if (i == 0) and (k == 0):

                        if z % 2 == 0:
                            t1x = (
                                xp * unit_x[0]
                                + yp * unit_y[0]
                                - (unit_x[0] + unit_y[0]) / 3
                                + unit_y[0] * length
                                + unit_x[0] * length
                            )
                            t1y = (
                                xp * unit_x[1]
                                + yp * unit_y[1]
                                - (unit_x[1] + unit_y[1]) / 3
                                + unit_y[1] * length
                                + unit_x[1] * length
                            )

                            t2x = t1x + unit_x[0]
                            t2y = t1y + unit_x[1]

                            t3x = t1x + unit_y[0]
                            t3y = t1y + unit_y[1]
                        else:
                            t1x = (
                                xp * unit_x[0]
                                + yp * unit_y[0]
                                + (unit_x[0] + unit_y[0]) / 3
                                + unit_y[0] * length
                                + unit_x[0] * length
                            )
                            t1y = (
                                xp * unit_x[1]
                                + yp * unit_y[1]
                                + (unit_x[1] + unit_y[1]) / 3
                                + unit_y[1] * length
                                + unit_x[1] * length
                            )

                            t2x = t1x - unit_x[0]
                            t2y = t1y - unit_x[1]

                            t3x = t1x - unit_y[0]
                            t3y = t1y - unit_y[1]

                        p = pat.Polygon(
                            xy=[(t1x, t1y), (t2x, t2y), (t3x, t3y)], fc=color, ec=color
                        )
                        ax.add_patch(p)

    p = pat.Polygon(
        xy=[(0, 0), (unit_y[0] * length, unit_y[1] * length), (0, unit_y[1] * length)],
        fc=(1, 1, 1),
        ec=(0, 0, 0),
    )
    ax.add_patch(p)

    p = pat.Polygon(
        xy=[
            (unit_x[0] * length, 0),
            ((unit_x[0] + unit_y[0]) * length, 0),
            ((unit_x[0] + unit_y[0]) * length, (unit_x[1] + unit_y[1]) * length),
        ],
        fc=(1, 1, 1),
        ec=(0, 0, 0),
    )
    ax.add_patch(p)

    ax.set_xlim(0, (unit_x[0] + unit_y[0]) * length)
    ax.set_ylim(0, (unit_x[1] + unit_y[1]) * length)
    ax.set_aspect("equal")

    return ax


def hist_formation(bx, n):
    numbers = layers_each[n]
    lay = []
    for i in range(len(numbers)):
        lay.append(i + 1)

    left: List[int] = list(range(numbers))

    bx.barh(left, numbers)
    bx.set_yticks(left)
    bx.set_yticklabels(lay)
    bx.set_xlabel("Coverage")
    bx.set_ylabel("layer number")

    bx.set_xlim(0, NBL)
    return bx


def figure_draw(n):
    ax = fig.add_subplot(211)
    ax.cla()
    ax = fig.add_subplot(211)
    kx = figure_formation(ax, n)
    # hist
    bx = fig.add_subplot(212)
    bx.cla()
    bx = fig.add_subplot(212)
    jx = hist_formation(bx, n)


def show_pictures():
    global canvas, imnum, fig
    root_p = tk.Toplevel()
    root_p.geometry("650x630")

    def p1_clicked():
        global imnum, cov_rec, canvas, fig
        imnum = imnum - 1
        if imnum < 0:
            imnum = len(cov_rec) - 1
        figure_draw(imnum)
        canvas.draw()
        canvas.get_tk_widget().pack()
        text_p8["text"] = str(int(t_rec[imnum])) + " s"
        text_p4["text"] = str(cov_rec[imnum])
        text_p6["text"] = str(imnum + 1) + "/" + str(len(cov_rec))

    def p2_clicked():
        global imnum, cov_rec, canvas, fig
        imnum = imnum + 1
        if imnum >= len(cov_rec):
            imnum = 0
        figure_draw(imnum)
        canvas.draw()
        canvas.get_tk_widget().pack()
        text_p8["text"] = str(int(t_rec[imnum])) + " s"
        text_p4["text"] = str(cov_rec[imnum])
        text_p6["text"] = str(imnum + 1) + "/" + str(len(cov_rec))

    button_p1 = tk.Button(
        root_p, text="Previous", command=p1_clicked, height=2, width=25
    )
    button_p1.place(x=30, y=560)
    button_p2 = tk.Button(root_p, text="Next", command=p2_clicked, height=2, width=25)
    button_p2.place(x=230, y=560)
    text_p7 = tk.Label(root_p, text="Time: ", font=("", 16))
    text_p7.place(x=30, y=500)
    text_p8 = tk.Label(root_p, text=str(cov_rec[0]), font=("", 16))
    text_p8.place(x=130, y=500)
    text_p3 = tk.Label(root_p, text="Coverage: ", font=("", 16))
    text_p3.place(x=230, y=500)
    text_p4 = tk.Label(root_p, text=str(cov_rec[0]), font=("", 16))
    text_p4.place(x=350, y=500)
    text_p6 = tk.Label(root_p, text="1/" + str(len(cov_rec)), font=("", 16))
    text_p6.place(x=450, y=500)
    imnum = len(cov_rec) - 1
    fig = plt.Figure()
    text_p8["text"] = str(int(t_rec[imnum])) + " s"
    text_p4["text"] = str(cov_rec[imnum])
    text_p6["text"] = str(imnum + 1) + "/" + str(len(cov_rec))
    figure_draw(imnum)
    canvas = FigureCanvasTkAgg(fig, master=root_p)
    canvas.draw()
    canvas.get_tk_widget().pack()

    def button_closep_clicked():
        plt.close("all")
        root_p.destroy()

    button_closep = tk.Button(
        root_p, text="Close", command=button_closep_clicked, height=2, width=25
    )
    button_closep.place(x=430, y=560)
    root_p.mainloop()


def cal_start():
    global lattice, atom_set, bonds, event, mod, rel_time, tot
    global a_set_rec, cov_rec, t_rec, images, d_rate
    start = time.time()
    pbval = 0
    pb.configure(value=pbval)
    text_count["text"] = "Started"
    text_coverage["text"] = "0 ML"
    root.update()
    # lists for recording
    a_set_rec = []
    cov_rec = []
    t_rec = []
    images = []
    # form lattice
    lattice_form()
    params()
    # tot: sum of rate constants. First, only deposition
    tot = d_rate
    # evaporation start
    rel_time = 0
    # first deposition
    global n_atoms
    n_atoms = 0
    deposition()
    text_count["text"] = "0 %"
    root.update()
    time_progress()
    # record first deposition
    rec_atoms()
    # update events and rates
    update_events()
    # second atomdepsoition
    deposition()
    time_progress()
    update_events()
    # repetition
    kMC()
    # record final structure
    text_count["text"] = "Saving"
    rec_atoms()
    # make poscar
    rec_pos()
    # record figures
    layer_sum()
    figure_draw_rec()
    coverage_color()
    hist_rec()
    global minute, second
    elapsed_time = time.time() - start
    minute = math.floor(elapsed_time / 60)
    second = int(elapsed_time % 60)
    text_count["text"] = str(minute) + " min " + str(second) + " sec"
    text_count["text"] = "PPT forming"
    ppt_form()
    text_count["text"] = "Finished"
    root.update()
    # calculation end
    # Show results
    show_pictures()


class App(Window):
    def __init__(self, master):
        super().__init__(master)


if __name__ == "__main__":
    unit_x: List[float] = [1, 0, 0]
    unit_y: List[float] = [0.5, 0.866, 0]
    unit_z: List[float] = [0, 0, 1]
    application = tk.Tk()
    app = App(application)
    app.run()
